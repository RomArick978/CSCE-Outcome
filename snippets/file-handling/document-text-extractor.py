"""
SNIPPET: Document Text Extractor
CATEGORY: File Handling
LANGUAGE: Python (FastAPI)
STATUS: Ready

DESCRIPTION:
    Extracts text content from uploaded documents for AI analysis.
    Supports PDF, DOCX, DOC, XLSX/XLS, PPTX, CSV, JSON, XML, TXT, and Markdown.
    Works with FastAPI UploadFile or file paths.

DEPENDENCIES:
    pip install pypdf python-docx openpyxl python-pptx

USAGE:
    from document_text_extractor import extract_text, extract_text_from_path

    # With FastAPI UploadFile
    @app.post("/extract")
    async def extract(file: UploadFile):
        text = await extract_text(file)
        return {"text": text, "filename": file.filename}

    # With file path
    text = extract_text_from_path("/path/to/file.pdf", "application/pdf", "report.pdf")

    # Multi-file
    @app.post("/extract-many")
    async def extract_many(files: list[UploadFile]):
        results = []
        for f in files:
            text = await extract_text(f)
            results.append({"name": f.filename, "text": text})
        return {"documents": results}

RELATED:
    - snippets/llm/bayer-llm-client.py (send extracted text to LLM)
    - snippets/file-handling/upload-handler.py (general file uploads)
"""

import csv
import io
import json
import os
import tempfile
from pathlib import Path

from pypdf import PdfReader
import docx
import openpyxl
from pptx import Presentation

# ---------------------------------------------------------------------------
# Supported file types (for frontend validation)
# ---------------------------------------------------------------------------

SUPPORTED_TYPES = {
    "application/pdf": "PDF",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "DOCX",
    "application/msword": "DOC",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "XLSX",
    "application/vnd.ms-excel": "XLS",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "PPTX",
    "text/plain": "TXT",
    "text/csv": "CSV",
    "text/markdown": "Markdown",
    "application/json": "JSON",
    "application/xml": "XML",
}

SUPPORTED_EXTENSIONS = ".pdf,.docx,.doc,.xlsx,.xls,.pptx,.csv,.json,.xml,.txt,.md,.yaml,.yml"

# ---------------------------------------------------------------------------
# Core: Extract text from file path
# ---------------------------------------------------------------------------

def extract_text_from_path(file_path: str, mimetype: str, filename: str) -> str:
    """Extract text content from a file on disk."""
    ext = Path(filename or "").suffix.lower().lstrip(".")

    # PDF
    if mimetype == "application/pdf" or ext == "pdf":
        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    # DOCX
    if (
        mimetype == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        or ext == "docx"
    ):
        doc = docx.Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs)

    # DOC (legacy) — best-effort
    if mimetype == "application/msword" or ext == "doc":
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception:
            return "[Could not extract text from .doc file]"

    # Excel (XLSX, XLS)
    if (
        mimetype in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        )
        or ext in ("xlsx", "xls")
    ):
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        sheets = []
        for name in wb.sheetnames:
            ws = wb[name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append(",".join(str(c) if c is not None else "" for c in row))
            sheets.append(f"[Sheet: {name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets)

    # PowerPoint (PPTX)
    if (
        mimetype == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or ext == "pptx"
    ):
        try:
            prs = Presentation(file_path)
            texts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_texts.append(shape.text)
                if slide_texts:
                    texts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
            return "\n\n".join(texts) or "[No text content found in PPTX]"
        except Exception:
            return "[Could not extract text from PPTX file]"

    # CSV
    if mimetype in ("text/csv", "application/csv") or ext == "csv":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    # Text-based files
    if (
        mimetype.startswith("text/")
        or mimetype in ("application/json", "application/xml")
        or ext in ("csv", "json", "xml", "txt", "md", "yaml", "yml", "toml")
    ):
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    # Fallback
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        return "[Binary file - content could not be extracted]"


# ---------------------------------------------------------------------------
# FastAPI UploadFile helper
# ---------------------------------------------------------------------------

async def extract_text(upload_file) -> str:
    """
    Extract text from a FastAPI UploadFile.
    Writes to a temp file, extracts, then cleans up.
    """
    content = await upload_file.read()
    suffix = Path(upload_file.filename or "").suffix

    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    try:
        return extract_text_from_path(
            tmp_path,
            upload_file.content_type or "",
            upload_file.filename or "",
        )
    finally:
        os.unlink(tmp_path)